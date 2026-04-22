"""
build_pitch_deck.py
Generates hearth_pitch_deck.pptx  -- light warm theme matching the frontend.
Run:  python build_pitch_deck.py
"""
import sys
if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Palette  (exact frontend CSS values)
# ---------------------------------------------------------------------------
BG_VOID    = RGBColor(0xFD, 0xF9, 0xF4)   # --bg-void   slide fill
BG_DEEP    = RGBColor(0xFA, 0xF5, 0xEE)   # --bg-deep
BG_SURFACE = RGBColor(0xF2, 0xE9, 0xDC)   # --bg-surface
BG_CARD    = RGBColor(0xFF, 0xFF, 0xFF)   # --bg-card
BG_CARD2   = RGBColor(0xFA, 0xF5, 0xEF)   # --bg-card-2
AMBER      = RGBColor(0xB8, 0x74, 0x30)   # --amber
GOLD       = RGBColor(0x9A, 0x5E, 0x14)   # --gold
FLAME      = RGBColor(0xD0, 0x4C, 0x14)   # --flame
INK        = RGBColor(0x26, 0x15, 0x08)   # --cream (dark text)
INK_DIM    = RGBColor(0x6A, 0x42, 0x20)   # --cream-dim
MUTED      = RGBColor(0x9E, 0x72, 0x48)   # --muted
BORDER     = RGBColor(0xE8, 0xD8, 0xBE)   # border colour
TEAL       = RGBColor(0x2A, 0x7A, 0x6E)   # contrast accent
TEAL_LT    = RGBColor(0xD4, 0xED, 0xEB)   # teal light
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------

def rect(slide, x, y, w, h, color, line_color=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def oval(slide, cx, cy, r, color):
    s = slide.shapes.add_shape(9, cx-r, cy-r, r*2, r*2)
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def rnd(slide, x, y, w, h, color, line_color=None, line_w=Pt(0)):
    """Rounded rectangle (MSO shape 5)."""
    s = slide.shapes.add_shape(5, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def tb(slide, text, x, y, w, h,
       size=14, bold=False, italic=False, color=INK,
       align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame; tf.word_wrap = wrap
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = font
    return box

def mlb(slide, lines, x, y, w, h, default_align=PP_ALIGN.LEFT, font="Calibri"):
    """Multi-line textbox. lines = list of (text, size, bold, italic, color, align?)"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        text, size, bold, italic, color = item[:5]
        align = item[5] if len(item) > 5 else default_align
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = color
        r.font.name = font
    return box

def divider(slide, x, y, w, color=AMBER, h=Inches(0.04)):
    rect(slide, x, y, w, h, color)

def badge(slide, x, y, w, h, text, bg, fg=WHITE, size=11):
    rnd(slide, x, y, w, h, bg)
    tb(slide, text, x, y, w, h, size=size, bold=True, color=fg,
       align=PP_ALIGN.CENTER)

def nav_strip(slide, title, subtitle="", bg=BG_SURFACE):
    rect(slide, 0, 0, W, Inches(1.25), bg)
    divider(slide, 0, Inches(1.22), W, AMBER, h=Inches(0.045))
    tb(slide, title,
       Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.7),
       size=30, bold=True, color=INK, font="Georgia")
    if subtitle:
        tb(slide, subtitle,
           Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.38),
           size=12, italic=True, color=MUTED)

def footer(slide, text="Hearth AI  |  BUAN 6V99  |  Spring 2026  |  Nikhil Gajula"):
    rect(slide, 0, Inches(7.18), W, Inches(0.32), BG_SURFACE)
    divider(slide, 0, Inches(7.17), W, BORDER, h=Inches(0.02))
    tb(slide, text, Inches(0.4), Inches(7.2), Inches(12.5), Inches(0.25),
       size=8, color=MUTED, align=PP_ALIGN.CENTER)

def stat_tile(slide, x, y, w, h, number, label, sub,
              num_color=AMBER, bg=BG_CARD, border=BORDER):
    rnd(slide, x, y, w, h, bg, border, Pt(1.5))
    tb(slide, number,
       x+Inches(0.1), y+Inches(0.1), w-Inches(0.2), Inches(0.9),
       size=40, bold=True, color=num_color, align=PP_ALIGN.CENTER, font="Georgia")
    tb(slide, label,
       x+Inches(0.1), y+Inches(0.95), w-Inches(0.2), Inches(0.5),
       size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    tb(slide, sub,
       x+Inches(0.1), y+Inches(1.42), w-Inches(0.2), Inches(0.5),
       size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

def icon_card(slide, x, y, w, h, icon, title, body, accent=AMBER):
    rnd(slide, x, y, w, h, BG_CARD, BORDER, Pt(1))
    rect(slide, x, y, w, Inches(0.07), accent)          # top accent bar
    tb(slide, icon, x, y+Inches(0.12), w, Inches(0.7),
       size=26, align=PP_ALIGN.CENTER)
    tb(slide, title, x+Inches(0.1), y+Inches(0.82), w-Inches(0.2), Inches(0.45),
       size=13, bold=True, color=INK, align=PP_ALIGN.CENTER)
    divider(slide, x+Inches(0.25), y+Inches(1.24), w-Inches(0.5), BORDER)
    tb(slide, body, x+Inches(0.12), y+Inches(1.32), w-Inches(0.24), h-Inches(1.4),
       size=10, color=INK_DIM, wrap=True)

def bullet_list(slide, items, x, y, w, size=11, color=INK_DIM, gap=Inches(0.44)):
    for i, item in enumerate(items):
        tb(slide, "  " + item, x, y + i*gap, w, gap,
           size=size, color=color)

# ---------------------------------------------------------------------------
# SLIDE 1  --  HERO / TITLE
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)

# Full warm background
rect(s, 0, 0, W, H, BG_VOID)

# Left amber accent bar
rect(s, 0, 0, Inches(0.55), H, AMBER)
rect(s, 0, 0, Inches(0.22), H, GOLD)

# Decorative warm circles top-right
for cx, cy, r, clr in [
    (Inches(12.2), Inches(1.1), Inches(2.0), RGBColor(0xF2,0xE0,0xC4)),
    (Inches(13.0), Inches(2.3), Inches(1.4), RGBColor(0xEC,0xD5,0xB0)),
    (Inches(11.5), Inches(2.5), Inches(1.0), RGBColor(0xF5,0xE8,0xD4)),
    (Inches(12.6), Inches(0.4), Inches(0.7), RGBColor(0xE8,0xCC,0x9C)),
]:
    oval(s, cx, cy, r, clr)

# Flame dot accent
oval(s, Inches(0.85), Inches(1.55), Inches(0.32), FLAME)
oval(s, Inches(0.85), Inches(2.35), Inches(0.22), AMBER)
oval(s, Inches(0.85), Inches(3.0),  Inches(0.15), GOLD)

# Tag line
tb(s, "BUAN 6V99  |  AI Agentic Systems  |  Spring 2026",
   Inches(0.9), Inches(0.38), Inches(9), Inches(0.42),
   size=11, italic=True, color=MUTED)

# Main title
mlb(s, [
    ("Hearth", 72, True, False, AMBER, PP_ALIGN.LEFT),
], Inches(0.9), Inches(0.95), Inches(8.5), Inches(1.5), font="Georgia")

# Flame emoji
tb(s, "flame", Inches(5.2), Inches(0.92), Inches(1.5), Inches(1.2),
   size=62, align=PP_ALIGN.LEFT)
# use text character since emoji may not render cross-platform in pptx reliably --
# replace with a unicode flame symbol
tb(s, "\U0001F525", Inches(5.0), Inches(0.85), Inches(1.5), Inches(1.3),
   size=58, align=PP_ALIGN.LEFT, color=FLAME)

# Subtitle
tb(s, "An AI-Powered Memory Companion for Alzheimer's Caregivers",
   Inches(0.9), Inches(2.15), Inches(8.8), Inches(0.75),
   size=21, italic=True, color=INK_DIM, font="Calibri Light")

divider(s, Inches(0.9), Inches(2.95), Inches(6.8), AMBER)

# Description block
mlb(s, [
    ("When a lucid moment happens -- be ready.", 15, False, True, INK),
    ("", 7, False, False, INK),
    ("Hearth uses OpenAI GPT-4o to generate a personalized care package in seconds:", 13, False, False, INK_DIM),
    ("a family letter, a voice message, a photo story, and a caregiver dialogue guide.", 13, False, False, INK_DIM),
    ("", 7, False, False, INK),
    ("Built for caregivers. Grounded in real family memories. Ready before the window closes.", 12, False, True, MUTED),
], Inches(0.9), Inches(3.1), Inches(8.8), Inches(2.2))

# Right info panel
rnd(s, Inches(9.8), Inches(1.6), Inches(3.1), Inches(5.1), BG_SURFACE, BORDER, Pt(1.5))
tb(s, "At a Glance", Inches(9.95), Inches(1.75), Inches(2.8), Inches(0.45),
   size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
divider(s, Inches(10.1), Inches(2.18), Inches(2.5), AMBER, h=Inches(0.03))

glance = [
    ("Platform", "Web App (FastAPI + JS SPA)"),
    ("AI Engine", "OpenAI GPT-4o"),
    ("Artifacts", "4 per session"),
    ("Time to generate", "< 60 seconds"),
    ("Demo mode", "No API key needed"),
    ("Data storage", "Local JSON files"),
    ("Voice support", "ElevenLabs TTS"),
    ("Users", "Caregivers + Families"),
]
for i, (k, v) in enumerate(glance):
    gy = Inches(2.32) + i * Inches(0.48)
    tb(s, k + ":", Inches(10.0), gy, Inches(1.15), Inches(0.42),
       size=10, bold=True, color=INK_DIM)
    tb(s, v, Inches(11.12), gy, Inches(1.65), Inches(0.42),
       size=10, color=INK)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 2  --  THE PROBLEM: ALZHEIMER'S BY THE NUMBERS
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Alzheimer's Disease -- By the Numbers",
          "The scale of the challenge and why families are caught unprepared")

# 5 stat tiles
tile_w = Inches(2.3)
tile_h = Inches(1.95)
gap    = Inches(0.26)
row_y  = Inches(1.42)
sx0    = Inches(0.38)

stats = [
    ("55M+",  "People living\nwith dementia",        "Worldwide (WHO 2024)",   FLAME),
    ("6.9M",  "Americans with\nAlzheimer's disease",  "Expected 13M by 2050",  AMBER),
    ("#1",    "Most common\nform of dementia",         "60-70% of all cases",   GOLD),
    ("$360B", "Annual care cost\nin the United States","Lost productivity incl.",TEAL),
    ("11M+",  "Unpaid U.S.\ncaregivers",               "18.4 billion hrs/year", INK_DIM),
]
for i, (num, lbl, sub, nc) in enumerate(stats):
    stat_tile(s, sx0 + i*(tile_w+gap), row_y, tile_w, tile_h, num, lbl, sub, num_color=nc)

# Separator
divider(s, Inches(0.38), Inches(3.5), Inches(12.55), BORDER)

# Stage distribution visual (horizontal bars)
tb(s, "U.S. Alzheimer's Stage Distribution (estimated)", Inches(0.38), Inches(3.6),
   Inches(7.0), Inches(0.38), size=12, bold=True, color=INK_DIM)

stages = [
    ("Mild / Early Stage",    0.35, TEAL,  "35%  ~2.4 Million people"),
    ("Moderate Stage",        0.40, AMBER, "40%  ~2.8 Million people"),
    ("Severe / Late Stage",   0.25, FLAME, "25%  ~1.7 Million people"),
]
bar_total_w = Inches(6.8)
bh = Inches(0.38)
for j, (lbl, pct, clr, note) in enumerate(stages):
    by = Inches(4.05) + j*Inches(0.62)
    tb(s, lbl, Inches(0.38), by, Inches(2.1), bh,
       size=11, bold=True, color=INK_DIM)
    bx = Inches(2.55)
    rect(s, bx, by+Inches(0.04), bar_total_w, bh-Inches(0.08), BG_SURFACE)
    rect(s, bx, by+Inches(0.04), int(bar_total_w*pct), bh-Inches(0.08), clr)
    tb(s, note, bx+Inches(0.08), by, bar_total_w, bh,
       size=10, bold=True, color=WHITE)

# Right panel: quick facts
rnd(s, Inches(9.6), Inches(3.5), Inches(3.35), Inches(3.5), BG_CARD, BORDER, Pt(1))
rect(s, Inches(9.6), Inches(3.5), Inches(3.35), Inches(0.06), AMBER)
tb(s, "Key Facts", Inches(9.75), Inches(3.62), Inches(3.0), Inches(0.42),
   size=13, bold=True, color=GOLD)
facts = [
    "New diagnosis every 3 seconds (global)",
    "Most common in adults over 65",
    "Average survival 4-8 years post-diagnosis",
    "Two-thirds of patients are women",
    "Cases will nearly triple by 2050",
    "No cure exists -- care is the intervention",
]
for i, f in enumerate(facts):
    tb(s, f"▸  {f}", Inches(9.78), Inches(4.1)+i*Inches(0.47),
       Inches(3.05), Inches(0.44), size=10, color=INK_DIM)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 3  --  LUCID MOMENTS: THE WINDOW
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Lucid Moments -- The Window of Opportunity",
          "Unexpected, brief, precious -- and almost always wasted without preparation")

# Left explanation panel
rnd(s, Inches(0.35), Inches(1.42), Inches(5.6), Inches(5.5), BG_SURFACE, BORDER, Pt(1.5))
rect(s, Inches(0.35), Inches(1.42), Inches(5.6), Inches(0.07), AMBER)
tb(s, "What is a Lucid Moment?",
   Inches(0.5), Inches(1.55), Inches(5.3), Inches(0.55),
   size=16, bold=True, color=INK, font="Georgia")

lucid_body = (
    "A lucid moment is a sudden, temporary period of mental clarity "
    "in a person with Alzheimer's or dementia -- when they recognize "
    "loved ones, recall memories, and communicate meaningfully.\n\n"
    "These windows can last anywhere from a few minutes to several "
    "hours. They are completely unpredictable. They are deeply "
    "precious. And most families are caught entirely unprepared."
)
tb(s, lucid_body, Inches(0.5), Inches(2.15), Inches(5.25), Inches(2.2),
   size=11.5, color=INK_DIM, wrap=True)

divider(s, Inches(0.65), Inches(4.4), Inches(5.1), AMBER)

# Lucid stats
lucid_stats = [
    ("73%",   "of caregivers say they felt unprepared\nwhen a lucid moment happened",   FLAME),
    ("15-45", "average minutes of clarity before\nconfusion returns",                    AMBER),
    ("2-4x",  "more lucid moments occur in the\nearly morning than other times",         GOLD),
]
for i, (n, l, c) in enumerate(lucid_stats):
    sy = Inches(4.55) + i * Inches(0.64)
    tb(s, n, Inches(0.5), sy, Inches(1.1), Inches(0.58),
       size=22, bold=True, color=c, font="Georgia")
    tb(s, l, Inches(1.6), sy, Inches(4.1), Inches(0.58),
       size=10, color=INK_DIM)

# Right: what goes wrong
rx = Inches(6.3)
tb(s, "What Families Face Without Preparation",
   rx, Inches(1.42), Inches(6.7), Inches(0.5),
   size=15, bold=True, color=INK)
divider(s, rx, Inches(1.9), Inches(6.6), AMBER)

problems = [
    ("Scrambling for words",
     "No prepared message ready. Caregivers freeze. Families freeze. The silence is painful.",
     FLAME),
    ("No photos or stories at hand",
     "Meaningful memories live in family members' heads -- nowhere near the patient when needed.",
     AMBER),
    ("No shared context between caregivers",
     "Each staff member interacts differently. No one knows what soothes this specific patient.",
     GOLD),
    ("The window closes before connection happens",
     "By the time families gather their thoughts, the patient has drifted back into confusion.",
     TEAL),
]
for i, (title, desc, accent) in enumerate(problems):
    py = Inches(2.05) + i * Inches(1.22)
    rnd(s, rx, py, Inches(6.65), Inches(1.1), BG_CARD, BORDER, Pt(1))
    rect(s, rx, py, Inches(0.06), Inches(1.1), accent)
    tb(s, title, rx+Inches(0.2), py+Inches(0.08), Inches(6.3), Inches(0.4),
       size=12, bold=True, color=INK)
    tb(s, desc, rx+Inches(0.2), py+Inches(0.48), Inches(6.3), Inches(0.55),
       size=10.5, color=INK_DIM)

# Closing banner
rnd(s, Inches(0.35), Inches(7.0), Inches(12.55), Inches(0.35), BG_SURFACE, BORDER, Pt(1))
tb(s, "Hearth solves all four of these -- before the moment ever happens.",
   Inches(0.55), Inches(7.04), Inches(12.1), Inches(0.3),
   size=12, bold=True, italic=True, color=AMBER, align=PP_ALIGN.CENTER)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 4  --  HOW HEARTH WORKS (5-STEP PIPELINE)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "How Hearth Works",
          "From family memories to a care-ready package -- in under 60 seconds")

steps = [
    ("1", "\U0001F4F1", "Family\nRegisters", "Caregiver or family\nmember opens Hearth\nweb app and adds\npatient profile",      BG_SURFACE,  AMBER),
    ("2", "\U0001F4DD", "Memories\nSubmitted", "Family shares memories,\nupload photos, and\nrecord voice samples\nthrough the app", BG_SURFACE, GOLD),
    ("3", "⚡",     "Lucid Now\nTriggered", "When a lucid moment\nbegins, caregiver taps\n'Lucid Now' -- starting\nthe AI pipeline",   BG_SURFACE, FLAME),
    ("4", "\U0001F9E0", "GPT-4o\nGenerates", "OpenAI GPT-4o reads\nall memories and\ncreates 4 personalized\nartifacts in seconds", BG_SURFACE, TEAL),
    ("5", "\U0001F91D", "Caregiver\nDelivers", "Letter read aloud,\nphoto story shown,\nvoice played, guide\nfollowed -- together", BG_SURFACE, INK_DIM),
]

cw = Inches(2.2)
ch = Inches(4.8)
gap_s = Inches(0.3)
start_s = Inches(0.37)

for i, (num, icon, title, desc, bg, accent) in enumerate(steps):
    cx = start_s + i*(cw+gap_s)
    cy = Inches(1.42)

    rnd(s, cx, cy, cw, ch, bg, BORDER, Pt(1.5))
    rect(s, cx, cy, cw, Inches(0.07), accent)

    # Step badge
    oval(s, cx+Inches(0.36), cy+Inches(0.42), Inches(0.3), accent)
    tb(s, num, cx+Inches(0.1), cy+Inches(0.15), Inches(0.55), Inches(0.55),
       size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Icon
    tb(s, icon, cx, cy+Inches(0.7), cw, Inches(0.78),
       size=30, align=PP_ALIGN.CENTER, color=accent)

    # Title
    tb(s, title, cx+Inches(0.08), cy+Inches(1.45), cw-Inches(0.16), Inches(0.65),
       size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)

    divider(s, cx+Inches(0.3), cy+Inches(2.08), cw-Inches(0.6), accent, h=Inches(0.03))

    # Description
    tb(s, desc, cx+Inches(0.1), cy+Inches(2.18), cw-Inches(0.2), Inches(2.45),
       size=10.5, color=INK_DIM, align=PP_ALIGN.CENTER, wrap=True)

    # Arrow
    if i < len(steps)-1:
        ax = cx + cw + Inches(0.03)
        tb(s, "→", ax, cy+Inches(2.0), gap_s+Inches(0.05), Inches(0.55),
           size=24, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Bottom note
rnd(s, Inches(0.35), Inches(6.42), Inches(12.55), Inches(0.55), BG_SURFACE, BORDER, Pt(1))
tb(s, "All processing is server-side  |  OpenAI GPT-4o powers every artifact  |  Demo mode works without an API key",
   Inches(0.55), Inches(6.5), Inches(12.1), Inches(0.38),
   size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 5  --  THE FOUR ARTIFACTS
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "The Four Artifacts",
          "Personalized for every patient -- generated by GPT-4o from real family memories")

artifacts = [
    ("✉️", "Family Letter",
     "A ~200-word letter written in the family member's voice, using real names, places, and shared memories. "
     "Warm and simple -- read aloud by the caregiver to open the visit.",
     ["Written in family member's voice", "Real names & places used", "~200 words, print-ready PDF", "No medical language"],
     AMBER),
    ("\U0001F399️", "Voice Message",
     "A 60-80 word spoken script -- intimate and warm, like a voicemail from a loved one. "
     "Can be read aloud or recorded. Hearing a familiar voice anchors the patient in the moment.",
     ["Exactly 60-80 words", "Natural spoken language", "Specific memory references", "Warm, personal closing"],
     GOLD),
    ("\U0001F5BC️", "Photo Story",
     "A PDF with family photos and caregiver-read captions. Each photo gets a 2-3 sentence caption -- "
     "sensory, emotional, specific -- designed to be read aloud while showing the photo.",
     ["One page per photo", "18pt+ accessible font", "Caption tied to real memory", "Embeds actual uploaded photos"],
     FLAME),
    ("\U0001F4AC", "Dialogue Guide",
     "A full visit script for the caregiver: greeting, 5 memory prompts, photo sequence, distress "
     "de-escalation script, and gentle closing -- built from the patient's emotional anchors.",
     ["Opening & closing scripts", "5 tailored memory prompts", "Photo show sequence", "Distress response script"],
     TEAL),
]

cw2 = Inches(3.0)
ch2 = Inches(5.55)
gap2 = Inches(0.24)
sx2  = Inches(0.36)

for i, (icon, title, desc, bullets, accent) in enumerate(artifacts):
    cx = sx2 + i*(cw2+gap2)
    cy = Inches(1.42)

    rnd(s, cx, cy, cw2, ch2, BG_CARD, BORDER, Pt(1.5))
    rect(s, cx, cy, cw2, Inches(0.07), accent)

    oval(s, cx+cw2/2, cy+Inches(0.7), Inches(0.46), BG_SURFACE)
    tb(s, icon, cx, cy+Inches(0.14), cw2, Inches(0.75),
       size=26, align=PP_ALIGN.CENTER, color=accent)

    tb(s, title, cx+Inches(0.1), cy+Inches(0.9), cw2-Inches(0.2), Inches(0.48),
       size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)

    divider(s, cx+Inches(0.25), cy+Inches(1.36), cw2-Inches(0.5), accent)

    tb(s, desc, cx+Inches(0.12), cy+Inches(1.48), cw2-Inches(0.24), Inches(1.85),
       size=10, color=INK_DIM, wrap=True)

    divider(s, cx+Inches(0.25), cy+Inches(3.3), cw2-Inches(0.5), BORDER)

    for j, b in enumerate(bullets):
        tb(s, f"▸  {b}", cx+Inches(0.14), cy+Inches(3.42)+j*Inches(0.5),
           cw2-Inches(0.22), Inches(0.46), size=10, color=INK_DIM)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 6  --  CARETAKER AI CHAT
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Caretaker AI Chat -- Ask Hearth Anything",
          "Context-aware guidance powered by GPT-4o, scoped strictly to dementia care")

# Left: feature overview
rnd(s, Inches(0.35), Inches(1.42), Inches(5.9), Inches(5.55), BG_SURFACE, BORDER, Pt(1.5))
rect(s, Inches(0.35), Inches(1.42), Inches(5.9), Inches(0.07), TEAL)
tb(s, "What the Chat Interface Does",
   Inches(0.5), Inches(1.56), Inches(5.6), Inches(0.48),
   size=15, bold=True, color=INK)

chat_desc = (
    "Caregivers can have a real-time conversation with Hearth's AI -- asking questions "
    "about care strategies, communication techniques, or patient-specific approaches.\n\n"
    "When a patient is selected from the dropdown, GPT-4o receives the patient's full profile "
    "(stage, emotional anchors, calming topics, memory count) as context -- making every "
    "answer specific to that person, not generic advice."
)
tb(s, chat_desc, Inches(0.5), Inches(2.1), Inches(5.6), Inches(2.0),
   size=11, color=INK_DIM, wrap=True)

divider(s, Inches(0.65), Inches(4.15), Inches(5.4), AMBER)

features = [
    "Real-time token streaming (no wait for full response)",
    "Patient context injected automatically into system prompt",
    "Conversation history maintained across turns (last 20)",
    "Strictly scoped: off-topic questions politely declined by GPT-4o",
    "Clear chat button resets conversation history",
    "Works without patient selected (general dementia guidance)",
]
for i, f in enumerate(features):
    tb(s, f"▸  {f}", Inches(0.52), Inches(4.32)+i*Inches(0.44),
       Inches(5.55), Inches(0.42), size=10.5, color=INK_DIM)

# Right: mock chat bubbles
rx = Inches(6.6)
tb(s, "Sample Interaction", rx, Inches(1.42), Inches(6.4), Inches(0.45),
   size=14, bold=True, color=INK)
divider(s, rx, Inches(1.85), Inches(6.3), AMBER)

chat_msgs = [
    ("Caregiver", "Dorothy seems agitated today. She keeps asking for Walter. What should I do?", False),
    ("Hearth AI", "Walter was Dorothy's husband of 51 years -- this name is deeply important to her. "
                  "Try gently saying: 'Walter loved you very much. He would want you to feel peaceful right now.' "
                  "Then redirect to her calming topics -- hymns or the smell of rain. Playing a familiar hymn "
                  "on low volume often helps ground her.", True),
    ("Caregiver", "Should I tell her Walter has passed away?", False),
    ("Hearth AI", "In dementia care, it is rarely helpful to repeat news of a loss -- it can cause fresh grief "
                  "each time. Instead, validate her feelings and redirect warmly toward comfort and memory.", True),
]
msg_y = Inches(2.0)
for (who, msg, is_ai) in chat_msgs:
    bg_c  = TEAL_LT if is_ai else RGBColor(0xFF, 0xF0, 0xE0)
    ac    = TEAL if is_ai else AMBER
    align = PP_ALIGN.LEFT
    mx    = rx if is_ai else rx + Inches(1.5)
    mw    = Inches(5.7) if is_ai else Inches(4.5)
    # estimate height
    n_lines = max(1, len(msg)//68 + 1)
    mh = Inches(0.3) + n_lines * Inches(0.3)
    rnd(s, mx, msg_y, mw, mh+Inches(0.35), bg_c, ac, Pt(0.75))
    tb(s, who.upper(), mx+Inches(0.12), msg_y+Inches(0.06), mw-Inches(0.2), Inches(0.25),
       size=8, bold=True, color=ac)
    tb(s, msg, mx+Inches(0.12), msg_y+Inches(0.28), mw-Inches(0.2), mh,
       size=9.5, color=INK_DIM, wrap=True)
    msg_y += mh + Inches(0.5)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 7  --  TECHNOLOGY ARCHITECTURE
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Technology Architecture",
          "FastAPI  |  OpenAI GPT-4o  |  Vanilla JS SPA  |  fpdf2  |  ElevenLabs")

cols = [
    ("Frontend (SPA)", TEAL, [
        "Vanilla JavaScript -- no framework",
        "7 screens, zero page reloads",
        "Server-Sent Events for live\njob progress streaming",
        "Token-by-token chat UI",
        "Drag & drop file uploads",
        "Manual history stack\n(back/forward nav)",
        "CSS custom properties\n(light warm theme)",
    ]),
    ("Backend (FastAPI)", AMBER, [
        "Python 3.11+  |  FastAPI  |  Uvicorn",
        "REST API + SSE streaming endpoints",
        "ThreadPoolExecutor for\nnon-blocking AI calls",
        "In-memory job store with\nreal-time status updates",
        "Patient data: JSON files\n(profile + memories)",
        "CORS middleware for local dev",
        "Demo mode -- no API key needed",
    ]),
    ("AI & Output Layer", FLAME, [
        "OpenAI GPT-4o (artifact gen)",
        "GPT-4o streaming (chat)",
        "fpdf2 -- all PDF generation",
        "Pillow -- photo embedding\nwith BytesIO (no file lock)",
        "ElevenLabs TTS (optional)",
        "Per-patient context\nin every system prompt",
        "JSON output with\nmarkdown fence stripping",
    ]),
]

cw3 = Inches(4.0)
gap3 = Inches(0.24)
sx3  = Inches(0.35)

for i, (title, accent, items) in enumerate(cols):
    cx = sx3 + i*(cw3+gap3)
    cy = Inches(1.42)
    ch3 = Inches(5.6)

    rnd(s, cx, cy, cw3, ch3, BG_CARD, BORDER, Pt(1.5))
    rect(s, cx, cy, cw3, Inches(0.07), accent)

    tb(s, title, cx+Inches(0.15), cy+Inches(0.13), cw3-Inches(0.3), Inches(0.5),
       size=16, bold=True, color=accent)
    divider(s, cx+Inches(0.2), cy+Inches(0.62), cw3-Inches(0.4), BG_SURFACE)

    for j, item in enumerate(items):
        tb(s, f"▸  {item}",
           cx+Inches(0.18), cy+Inches(0.75)+j*Inches(0.66),
           cw3-Inches(0.3), Inches(0.62), size=10.5, color=INK_DIM)

# Data flow banner
rnd(s, Inches(0.35), Inches(7.12), Inches(12.55), Inches(0.28), BG_SURFACE, BORDER, Pt(1))
tb(s, "Data flow:  Browser  <->  FastAPI (REST + SSE)  <->  GPT-4o Agent  ->  PDF Generators  ->  File Store  ->  Browser",
   Inches(0.5), Inches(7.15), Inches(12.2), Inches(0.24),
   size=9, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 8  --  EVALUATION PIPELINE
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Evaluation Pipeline -- How We Measure Success",
          "Artifact quality, user experience, and clinical relevance across three dimensions")

# Dimension cards
dims = [
    ("Artifact Quality\nEvaluation", AMBER,
     "How good are the AI-generated outputs?",
     [
        "Expert panel review (dementia care specialists grade each artifact)",
        "Rubric: specificity, warmth, avoidance of medical language",
        "BLEU / ROUGE vs. human-written gold-standard letters",
        "Caregiver blind comparison: AI artifact vs. human-written",
        "Readability scoring: Flesch-Kincaid target 6th grade level",
     ]),
    ("User Experience\nEvaluation", TEAL,
     "Does the interface reduce caregiver stress?",
     [
        "System Usability Scale (SUS) survey after first use",
        "Time-to-generate measurement: target < 60 seconds",
        "Caregiver confidence score before/after using the package",
        "Task completion rate for memory submission workflow",
        "Net Promoter Score from caregiver cohort",
     ]),
    ("Clinical Impact\nEvaluation", FLAME,
     "Does the patient benefit?",
     [
        "Patient engagement duration during lucid moments (minutes)",
        "Caregiver-reported anxiety level of patient during visit",
        "Repeat use rate: do caregivers run Lucid Now again?",
        "Family satisfaction score post-session",
        "Longitudinal: mood baseline improvement over 30 days",
     ]),
]

cw4 = Inches(4.1)
gap4 = Inches(0.25)
sx4  = Inches(0.33)

for i, (title, accent, question, items) in enumerate(dims):
    cx = sx4 + i*(cw4+gap4)
    cy = Inches(1.42)
    ch4 = Inches(5.58)

    rnd(s, cx, cy, cw4, ch4, BG_CARD, BORDER, Pt(1.5))
    rect(s, cx, cy, cw4, Inches(0.07), accent)

    tb(s, title, cx+Inches(0.15), cy+Inches(0.13), cw4-Inches(0.3), Inches(0.65),
       size=15, bold=True, color=accent)
    tb(s, question, cx+Inches(0.15), cy+Inches(0.76), cw4-Inches(0.3), Inches(0.4),
       size=10.5, italic=True, color=MUTED)
    divider(s, cx+Inches(0.2), cy+Inches(1.14), cw4-Inches(0.4), BG_SURFACE)

    for j, item in enumerate(items):
        rnd(s, cx+Inches(0.15), cy+Inches(1.25)+j*Inches(0.83),
            cw4-Inches(0.3), Inches(0.75), BG_SURFACE, BORDER, Pt(0.5))
        tb(s, f"  {item}",
           cx+Inches(0.22), cy+Inches(1.32)+j*Inches(0.83),
           cw4-Inches(0.4), Inches(0.65), size=10, color=INK_DIM)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 9  --  TRL LEVEL & IMPLEMENTATION READINESS
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Technology Readiness Level (TRL)",
          "Where Hearth stands today and what each next stage requires")

# TRL scale visual
trl_data = [
    (1, "Basic Principles",      "Concept and core idea validated",                           False),
    (2, "Technology Concept",    "Use case defined, target users identified",                  False),
    (3, "Proof of Concept",      "Core functionality demonstrated in lab environment",         False),
    (4, "Lab Validation",        "Key components tested -- AI outputs, PDF generation, UI",    False),
    (5, "Relevant Environment",  "Hearth tested with simulated patient profiles and memories", True),   # current
    (6, "Demo Environment",      "Functional prototype tested with real caregivers in demos",  False),
    (7, "Pilot Deployment",      "Pilot at one memory care facility with live patients",       False),
    (8, "System Complete",       "Validated in operational clinical environment",              False),
    (9, "Operational",           "Full deployment, HIPAA compliance, care network integration",False),
]

bar_h_each = Inches(0.56)
bar_w_full  = Inches(7.8)
bx0  = Inches(5.2)
by0  = Inches(1.38)

for i, (lvl, name, desc, current) in enumerate(trl_data):
    by = by0 + i * bar_h_each
    fill_pct = lvl / 9.0
    bg_row = BG_SURFACE if i % 2 == 0 else BG_CARD
    rect(s, Inches(0.35), by, Inches(12.6), bar_h_each-Inches(0.04), bg_row)

    # TRL number badge
    badge_color = AMBER if current else (BG_SURFACE if lvl > 5 else GOLD)
    badge_text_color = WHITE if current or lvl <= 5 else MUTED
    oval(s, Inches(0.74), by+bar_h_each/2, Inches(0.26),
         FLAME if current else (AMBER if lvl <= 5 else BORDER))
    tb(s, str(lvl), Inches(0.5), by+Inches(0.1), Inches(0.55), bar_h_each-Inches(0.12),
       size=12, bold=True, color=WHITE if current or lvl<=5 else MUTED,
       align=PP_ALIGN.CENTER)

    # Name
    fw = Inches(2.0)
    tb(s, name, Inches(1.1), by+Inches(0.08), fw, bar_h_each-Inches(0.1),
       size=10.5, bold=(lvl <= 5), color=FLAME if current else (INK if lvl<=5 else MUTED))

    # Progress bar fill
    fill_w = int(bar_w_full * fill_pct)
    rect(s, bx0, by+Inches(0.14), bar_w_full, bar_h_each-Inches(0.32), BG_SURFACE)
    fill_c = FLAME if current else (AMBER if lvl<=5 else BORDER)
    rect(s, bx0, by+Inches(0.14), fill_w, bar_h_each-Inches(0.32), fill_c)

    # Desc
    tb(s, desc, Inches(3.15), by+Inches(0.08), Inches(2.0), bar_h_each-Inches(0.1),
       size=9.5, color=INK_DIM if lvl<=5 else MUTED)

    # "CURRENT" label
    if current:
        badge(s, Inches(12.25), by+Inches(0.1), Inches(0.65), Inches(0.35),
              "NOW", FLAME, WHITE, size=9)

# Legend
rnd(s, Inches(0.35), Inches(6.6), Inches(12.55), Inches(0.52), BG_SURFACE, BORDER, Pt(1))
for lbl, clr, x in [
    ("Completed", GOLD, Inches(0.6)), ("Current TRL-5", FLAME, Inches(2.2)), ("Not yet reached", BORDER, Inches(4.2))
]:
    oval(s, x+Inches(0.15), Inches(6.87), Inches(0.12), clr)
    tb(s, lbl, x+Inches(0.35), Inches(6.68), Inches(1.5), Inches(0.38),
       size=10, color=INK_DIM)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 10  --  BUSINESS VERSION & GO-TO-MARKET
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Business Implementation Plan",
          "From academic prototype to a sustainable, scalable care-technology product")

# Market size top strip
mkt_items = [
    ("$7.5B",  "U.S. dementia care\ntech market (2024)",   FLAME),
    ("16,000+","Memory care facilities\nin the United States", AMBER),
    ("11M+",   "Potential end users\n(unpaid caregivers)",   GOLD),
    ("$360B",  "Annual dementia\ncare cost (target pain)",   TEAL),
]
mw = Inches(2.95)
for i, (num, lbl, c) in enumerate(mkt_items):
    mx = Inches(0.35) + i*(mw+Inches(0.22))
    rnd(s, mx, Inches(1.42), mw, Inches(1.18), BG_SURFACE, BORDER, Pt(1))
    rect(s, mx, Inches(1.42), mw, Inches(0.06), c)
    tb(s, num, mx+Inches(0.1), Inches(1.5), mw-Inches(0.2), Inches(0.62),
       size=30, bold=True, color=c, align=PP_ALIGN.CENTER, font="Georgia")
    tb(s, lbl, mx+Inches(0.1), Inches(2.1), mw-Inches(0.2), Inches(0.44),
       size=10, color=INK_DIM, align=PP_ALIGN.CENTER)

divider(s, Inches(0.35), Inches(2.72), Inches(12.6), BORDER)

# Business model and GTM side by side
for col_x, col_w, title, accent, items in [
    (Inches(0.35), Inches(5.9), "Business Model (SaaS)", AMBER, [
        ("Subscription Tiers",
         "B2B: $199/month per facility (unlimited patients)\nB2C: $9.99/month per family (1 patient)"),
        ("Freemium Entry",
         "Free demo mode with pre-generated content -- converts families to paid on first lucid moment"),
        ("Enterprise",
         "Health system licensing, branded white-label builds, EHR integration modules"),
        ("Data Moat",
         "De-identified interaction patterns inform prompt refinement -- improving output quality over time"),
    ]),
    (Inches(6.6), Inches(5.9), "Go-To-Market Strategy", TEAL, [
        ("Phase 1: Direct to Families (0-6 months)",
         "Social media targeting family caregivers; Alzheimer's Association partnership; free tier drive sign-ups"),
        ("Phase 2: Memory Care Facilities (6-18 months)",
         "Pilot with 3-5 facilities; case study ROI data; staff training program; facility branding option"),
        ("Phase 3: Health System Integration (18-36 months)",
         "HIPAA compliance build-out; EHR API integration; insurance reimbursement pathway exploration"),
        ("Phase 4: International Expansion (36+ months)",
         "Localization for Spanish, Mandarin, German; WHO dementia partnership; multi-language GPT-4o prompting"),
    ]),
]:
    rnd(s, col_x, Inches(2.85), col_w, Inches(4.1), BG_CARD, BORDER, Pt(1.5))
    rect(s, col_x, Inches(2.85), col_w, Inches(0.07), accent)
    tb(s, title, col_x+Inches(0.15), Inches(2.98), col_w-Inches(0.3), Inches(0.45),
       size=14, bold=True, color=accent)
    divider(s, col_x+Inches(0.2), Inches(3.4), col_w-Inches(0.4), BG_SURFACE)

    iy = Inches(3.5)
    for (ititle, idesc) in items:
        tb(s, ititle, col_x+Inches(0.18), iy, col_w-Inches(0.3), Inches(0.35),
           size=11, bold=True, color=INK)
        tb(s, idesc, col_x+Inches(0.18), iy+Inches(0.33), col_w-Inches(0.3), Inches(0.5),
           size=9.5, color=INK_DIM, wrap=True)
        iy += Inches(0.92)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 11  --  PRODUCT BACKLOG & ROADMAP
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)
nav_strip(s, "Product Backlog & Roadmap",
          "Prioritized features, known limitations, and the path forward")

# Roadmap timeline top
phases = [
    ("v1.0 -- NOW",      "Academic Prototype",   FLAME,
     ["Core artifact generation", "5 demo patients", "Chat interface", "OpenAI integration"]),
    ("v1.5 -- Q3 2026",  "Beta Release",         AMBER,
     ["HIPAA data encryption", "Real ElevenLabs voice", "Mobile responsive UI", "Multi-language support"]),
    ("v2.0 -- Q1 2027",  "Production",           GOLD,
     ["EHR integration API", "Facility dashboard", "Analytics + reporting", "Family mobile app"]),
    ("v3.0 -- 2028",     "Platform",             TEAL,
     ["Health system partnerships", "Insurance reimbursement", "GPT fine-tuning on care data", "Wearable triggers"]),
]

pw = Inches(3.0)
pg = Inches(0.22)
px0 = Inches(0.35)

for i, (ver, phase_title, accent, feats) in enumerate(phases):
    px = px0 + i*(pw+pg)
    py = Inches(1.42)

    rnd(s, px, py, pw, Inches(2.8), BG_CARD, BORDER, Pt(1.5))
    rect(s, px, py, pw, Inches(0.07), accent)

    badge(s, px+Inches(0.1), py+Inches(0.14), pw-Inches(0.2), Inches(0.3),
          ver, accent, WHITE, size=10)
    tb(s, phase_title, px+Inches(0.12), py+Inches(0.52), pw-Inches(0.22), Inches(0.38),
       size=13, bold=True, color=INK)
    divider(s, px+Inches(0.25), py+Inches(0.88), pw-Inches(0.5), BG_SURFACE)
    for j, f in enumerate(feats):
        tb(s, f"✓  {f}", px+Inches(0.15), py+Inches(1.0)+j*Inches(0.43),
           pw-Inches(0.25), Inches(0.4), size=10, color=INK_DIM)

    # Arrow
    if i < len(phases)-1:
        ax = px + pw + Inches(0.02)
        tb(s, "→", ax, py+Inches(1.1), pg+Inches(0.08), Inches(0.5),
           size=22, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

divider(s, Inches(0.35), Inches(4.35), Inches(12.6), BORDER)

# Backlog + known risks two columns
for col_x, col_w, title, accent, items in [
    (Inches(0.35), Inches(6.05), "Prioritized Backlog (Next Sprint)", AMBER, [
        ("HIGH", "Add HIPAA-compliant encryption for patient data at rest"),
        ("HIGH", "Replace demo voice output with real ElevenLabs TTS audio"),
        ("HIGH", "Patient photo display in web UI (currently PDF only)"),
        ("MED",  "Caregiver session log -- track which artifacts were used"),
        ("MED",  "Add memory tagging (dates, people, places) for richer prompts"),
        ("MED",  "Automated artifact email delivery to on-call staff"),
        ("LOW",  "Dark mode toggle for night-shift caregivers"),
        ("LOW",  "Bulk patient CSV import for facility onboarding"),
    ]),
    (Inches(6.75), Inches(6.05), "Known Risks & Mitigations", FLAME, [
        ("RISK", "AI hallucination -- GPT-4o may invent details not in memories"),
        ("MIT",  "All prompts strictly reference submitted memories; human review step"),
        ("RISK", "No clinical validation -- outputs not reviewed by dementia specialists"),
        ("MIT",  "v1.5 includes expert review panel and rubric-based QA pipeline"),
        ("RISK", "Patient data privacy -- PII in JSON files is unencrypted"),
        ("MIT",  "v1.5 adds AES-256 encryption + access control per facility"),
        ("RISK", "Internet dependency -- no offline fallback for API calls"),
        ("MIT",  "Demo mode provides full offline functionality as fallback"),
    ]),
]:
    rnd(s, col_x, Inches(4.5), col_w, Inches(2.52), BG_CARD, BORDER, Pt(1.5))
    rect(s, col_x, Inches(4.5), col_w, Inches(0.07), accent)
    tb(s, title, col_x+Inches(0.15), Inches(4.6), col_w-Inches(0.3), Inches(0.42),
       size=13, bold=True, color=accent)
    divider(s, col_x+Inches(0.2), Inches(5.0), col_w-Inches(0.4), BG_SURFACE)
    for j, (tag, text) in enumerate(items):
        tx = col_x + Inches(0.18)
        ty = Inches(5.1) + j * Inches(0.3)
        tag_c = FLAME if tag in ("HIGH","RISK") else (AMBER if tag=="MED" else TEAL)
        badge(s, tx, ty+Inches(0.02), Inches(0.5), Inches(0.22), tag, tag_c, WHITE, size=7)
        tb(s, text, tx+Inches(0.58), ty, col_w-Inches(0.75), Inches(0.28),
           size=9.5, color=INK_DIM)

footer(s)


# ---------------------------------------------------------------------------
# SLIDE 12  --  CLOSING
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, BG_VOID)

# Warm ambient rings
for r_val, clr in [
    (Inches(4.2), RGBColor(0xF2,0xE9,0xDC)),
    (Inches(3.0), RGBColor(0xEC,0xDE,0xCA)),
    (Inches(1.8), RGBColor(0xE8,0xD4,0xB8)),
]:
    oval(s, W/2, H/2, r_val, clr)

# Left accent bar
rect(s, 0, 0, Inches(0.55), H, AMBER)
rect(s, 0, 0, Inches(0.22), H, GOLD)

# Flame icon centre
tb(s, "\U0001F525", Inches(0), Inches(0.95), W, Inches(1.4),
   size=72, align=PP_ALIGN.CENTER, color=FLAME)

# Title
tb(s, "Hearth",
   Inches(0), Inches(2.05), W, Inches(1.2),
   size=68, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Georgia")

divider(s, Inches(3.0), Inches(3.2), Inches(7.33), AMBER)

# Tagline
tb(s, "When a lucid moment happens -- be ready.",
   Inches(0.5), Inches(3.35), Inches(12.33), Inches(0.65),
   size=21, italic=True, color=INK_DIM, align=PP_ALIGN.CENTER, font="Calibri Light")

# Stack line
tb(s, "FastAPI  |  OpenAI GPT-4o  |  fpdf2  |  ElevenLabs  |  Vanilla JS",
   Inches(0.5), Inches(4.1), Inches(12.33), Inches(0.45),
   size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Creator line
tb(s, "Nikhil Gajula  |  BUAN 6V99 -- AI Agentic Systems  |  Spring 2026",
   Inches(0.5), Inches(4.62), Inches(12.33), Inches(0.42),
   size=12, color=MUTED, align=PP_ALIGN.CENTER)

divider(s, Inches(3.0), Inches(5.18), Inches(7.33), BORDER)

# Quote
tb(s, '"Technology in service of the most human moments."',
   Inches(1.5), Inches(5.35), Inches(10.33), Inches(0.75),
   size=16, italic=True, color=AMBER, align=PP_ALIGN.CENTER, font="Georgia")

# Bottom strip
rect(s, 0, Inches(7.1), W, Inches(0.4), BG_SURFACE)
divider(s, 0, Inches(7.08), W, BORDER, h=Inches(0.02))
tb(s, "All text in this presentation is fully editable. Built with python-pptx.",
   Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.28),
   size=8, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = Path(__file__).parent / "hearth_pitch_deck.pptx"
prs.save(str(out))
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
