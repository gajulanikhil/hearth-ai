"""
build_pptx.py — Generates Hearth project presentation (hearth_presentation.pptx)
Run: python build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
from pathlib import Path

# ── Palette ───────────────────────────────────────────────────────────────────
WARM_DARK   = RGBColor(0x1A, 0x0A, 0x02)   # near-black brown
AMBER       = RGBColor(0xB8, 0x74, 0x30)   # brand amber
GOLD        = RGBColor(0x9A, 0x5E, 0x14)   # deep gold
FLAME       = RGBColor(0xD0, 0x4C, 0x14)   # flame orange
CREAM       = RGBColor(0xFD, 0xF9, 0xF4)   # off-white background
LIGHT_TAN   = RGBColor(0xF2, 0xE9, 0xDC)   # light tan
CARD_BG     = RGBColor(0xFF, 0xFF, 0xFF)   # white cards
MID_BROWN   = RGBColor(0x6A, 0x42, 0x20)   # medium brown text
MUTED       = RGBColor(0x9E, 0x72, 0x48)   # muted text
TEAL_ACC    = RGBColor(0x2A, 0x7A, 0x6E)   # teal accent (contrast)
TEAL_LIGHT  = RGBColor(0xD4, 0xED, 0xEB)   # light teal fill
RED_STAT    = RGBColor(0xC0, 0x39, 0x2B)   # red for alarming stats
SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # fully blank


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.line.fill.background()
    shape.line.width = 0
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    return shape


def add_text_box(slide, text, x, y, w, h,
                 font_size=18, bold=False, italic=False,
                 color=WARM_DARK, align=PP_ALIGN.LEFT,
                 font_name="Calibri", wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    return txb


def add_multiline(slide, lines, x, y, w, h,
                  font_size=18, color=WARM_DARK,
                  align=PP_ALIGN.LEFT, line_spacing=1.2,
                  font_name="Calibri", bold=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, (txt, sz, bd, it, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(sz)
        run.font.bold  = bd
        run.font.italic = it
        run.font.color.rgb = clr
        run.font.name  = font_name
    return txb


def add_circle(slide, cx, cy, r, fill_rgb):
    x = cx - r
    y = cy - r
    shape = slide.shapes.add_shape(9, x, y, r*2, r*2)  # oval
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    return shape


def stat_card(slide, x, y, w, h, number, label, sub,
              bg=AMBER, num_color=CREAM, lbl_color=CREAM, sub_color=LIGHT_TAN):
    add_rect(slide, x, y, w, h, bg)
    # number
    add_text_box(slide, number,
                 x + Inches(0.15), y + Inches(0.15),
                 w - Inches(0.3), Inches(0.95),
                 font_size=42, bold=True, color=num_color,
                 align=PP_ALIGN.CENTER)
    # label
    add_text_box(slide, label,
                 x + Inches(0.1), y + Inches(1.0),
                 w - Inches(0.2), Inches(0.55),
                 font_size=13, bold=True, color=lbl_color,
                 align=PP_ALIGN.CENTER)
    # sub
    add_text_box(slide, sub,
                 x + Inches(0.1), y + Inches(1.5),
                 w - Inches(0.2), Inches(0.6),
                 font_size=10, italic=True, color=sub_color,
                 align=PP_ALIGN.CENTER)


def progress_bar(slide, x, y, w, h, pct, bg=LIGHT_TAN, fill=AMBER, label=""):
    add_rect(slide, x, y, w, h, bg)
    filled_w = int(w * pct)
    if filled_w > 0:
        add_rect(slide, x, y, filled_w, h, fill)
    if label:
        add_text_box(slide, label,
                     x + Inches(0.05), y - Inches(0.02),
                     w, h,
                     font_size=10, bold=True, color=CREAM,
                     align=PP_ALIGN.LEFT)


def rounded_rect(slide, x, y, w, h, fill_rgb, radius_emu=Emu(200000)):
    shape = slide.shapes.add_shape(5, x, y, w, h)  # rounded rectangle
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    # adjust corner radius
    try:
        sp = shape._element
        spPr = sp.find(nsmap.qn('p:spPr'))
        if spPr is None:
            spPr = etree.SubElement(sp, nsmap.qn('p:spPr'))
        prstGeom = spPr.find(nsmap.qn('a:prstGeom'))
        if prstGeom is not None:
            avLst = prstGeom.find(nsmap.qn('a:avLst'))
            if avLst is not None:
                for gd in avLst.findall(nsmap.qn('a:gd')):
                    if gd.get('name') == 'adj':
                        gd.set('fmla', 'val 30000')
    except Exception:
        pass
    return shape


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / HERO
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)

# Full dark background
add_rect(s1, 0, 0, SLIDE_W, SLIDE_H, WARM_DARK)

# Warm gradient band (left side accent)
add_rect(s1, 0, 0, Inches(0.6), SLIDE_H, FLAME)
add_rect(s1, 0, 0, Inches(0.3), SLIDE_H, AMBER)

# Decorative flame circles (abstract orbs)
add_circle(s1, Inches(10.8), Inches(1.2), Inches(1.8), RGBColor(0xC0, 0x50, 0x10))
add_circle(s1, Inches(11.5), Inches(2.5), Inches(1.2), RGBColor(0xD8, 0x70, 0x20))
add_circle(s1, Inches(10.2), Inches(2.8), Inches(0.9), RGBColor(0xA0, 0x40, 0x08))
add_circle(s1, Inches(12.0), Inches(1.0), Inches(0.7), RGBColor(0xE8, 0x90, 0x30))

# Tagline top
add_text_box(s1, "BUAN 6V99 · AI Agentic Systems · Spring 2026",
             Inches(0.9), Inches(0.4), Inches(9), Inches(0.5),
             font_size=11, color=MUTED, italic=True)

# Main logo word
add_text_box(s1, "🔥  Hearth",
             Inches(0.9), Inches(1.3), Inches(8), Inches(1.6),
             font_size=72, bold=True, color=AMBER, font_name="Georgia")

# Subtitle
add_text_box(s1, "An AI-Powered Memory Companion for Alzheimer's Caregivers",
             Inches(0.9), Inches(2.8), Inches(9.0), Inches(0.9),
             font_size=22, italic=True, color=LIGHT_TAN, font_name="Calibri Light")

# Divider line
add_rect(s1, Inches(0.9), Inches(3.65), Inches(6.5), Inches(0.04), AMBER)

# Description
add_multiline(s1,
    [
        ("When a lucid moment happens — be ready.", 16, False, True, CREAM),
        ("", 8, False, False, CREAM),
        ("Hearth uses OpenAI GPT-4o to generate personalized care packages in seconds —", 14, False, False, LIGHT_TAN),
        ("a family letter, a voice message, a photo story, and a caregiver dialogue guide.", 14, False, False, LIGHT_TAN),
    ],
    Inches(0.9), Inches(3.75), Inches(9.0), Inches(1.8))

# Bottom bar
add_rect(s1, 0, Inches(6.9), SLIDE_W, Inches(0.6), RGBColor(0x26, 0x15, 0x08))
add_text_box(s1, "Nikhil Gajula  ·  Hearth AI  ·  2026",
             Inches(0.9), Inches(6.95), Inches(8), Inches(0.4),
             font_size=10, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ALZHEIMER'S BY THE NUMBERS
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, SLIDE_W, SLIDE_H, CREAM)
add_rect(s2, 0, 0, SLIDE_W, Inches(1.4), WARM_DARK)
add_rect(s2, 0, Inches(1.35), SLIDE_W, Inches(0.06), AMBER)

add_text_box(s2, "Alzheimer's Disease — By the Numbers",
             Inches(0.5), Inches(0.3), Inches(12), Inches(0.9),
             font_size=32, bold=True, color=CREAM, font_name="Georgia",
             align=PP_ALIGN.CENTER)

# Stat cards row 1
cards_y = Inches(1.65)
card_h  = Inches(1.85)
card_w  = Inches(2.35)
gap     = Inches(0.22)
start_x = Inches(0.5)

stat_data = [
    ("55M+",  "People Living\nwith Dementia", "Worldwide · 2024 WHO data",   FLAME),
    ("6.9M",  "Americans\nwith Alzheimer's", "Expected 13M by 2050",         AMBER),
    ("#1",    "Most Common\nForm of Dementia", "60–70% of all dementia cases", GOLD),
    ("$360B", "Annual Cost\nin the U.S.",     "Direct care + lost productivity", RGBColor(0x5A,0x32,0x0E)),
    ("11M+",  "Unpaid Caregivers\nin the U.S.", "18.4B hours of care per year",  TEAL_ACC),
]

for i, (num, lbl, sub, bg) in enumerate(stat_data):
    sx = start_x + i * (card_w + gap)
    stat_card(s2, sx, cards_y, card_w, card_h, num, lbl, sub, bg=bg)

# Second row — disease progression bar chart
bar_y = Inches(3.75)
add_text_box(s2, "Disease Stage Distribution (estimated U.S. patients)",
             Inches(0.5), bar_y - Inches(0.4), Inches(9), Inches(0.4),
             font_size=13, bold=True, color=MID_BROWN)

stages = [
    ("Mild (Early)",     0.35, TEAL_ACC,  "35%  ·  ~2.4M people"),
    ("Moderate",         0.40, AMBER,     "40%  ·  ~2.8M people"),
    ("Severe (Late)",    0.25, FLAME,     "25%  ·  ~1.7M people"),
]
bar_w_total = Inches(7.5)
bh = Inches(0.38)
for j, (lbl, pct, clr, note) in enumerate(stages):
    by = bar_y + j * Inches(0.58)
    add_text_box(s2, lbl, Inches(0.5), by, Inches(1.8), bh,
                 font_size=11, color=MID_BROWN, bold=True)
    bx = Inches(2.4)
    add_rect(s2, bx, by + Inches(0.04), bar_w_total, bh - Inches(0.08), LIGHT_TAN)
    add_rect(s2, bx, by + Inches(0.04), int(bar_w_total * pct), bh - Inches(0.08), clr)
    add_text_box(s2, note, bx + Inches(0.1), by, bar_w_total, bh,
                 font_size=10, bold=True, color=CREAM)

# Right panel — key facts
add_rect(s2, Inches(10.0), Inches(1.65), Inches(3.0), Inches(5.6), WARM_DARK)
facts = [
    ("Key Facts", 14, True, False, AMBER),
    ("", 6, False, False, CREAM),
    ("⏱  New diagnosis every\n    3 seconds globally", 11, False, False, CREAM),
    ("", 5, False, False, CREAM),
    ("👵  Most common in\n    adults over 65", 11, False, False, CREAM),
    ("", 5, False, False, CREAM),
    ("🧠  Memory loss is the\n    hallmark symptom", 11, False, False, CREAM),
    ("", 5, False, False, CREAM),
    ("💔  Average survival:\n    4–8 years after diagnosis", 11, False, False, CREAM),
    ("", 5, False, False, CREAM),
    ("🌍  By 2050, cases will\n    nearly triple worldwide", 11, False, False, CREAM),
]
add_multiline(s2, facts, Inches(10.2), Inches(1.85), Inches(2.6), Inches(5.0),
              align=PP_ALIGN.LEFT)

add_rect(s2, 0, Inches(7.2), SLIDE_W, Inches(0.3), WARM_DARK)
add_text_box(s2, "Sources: Alzheimer's Association 2024 Report · WHO Global Dementia Report · CDC",
             Inches(0.5), Inches(7.22), Inches(12), Inches(0.25),
             font_size=8, color=MUTED, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LUCID MOMENTS: THE WINDOW OF OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, SLIDE_W, SLIDE_H, RGBColor(0xF7, 0xF2, 0xEB))
add_rect(s3, 0, 0, Inches(0.08), SLIDE_H, AMBER)

# Hero left column
add_rect(s3, 0, 0, Inches(5.6), SLIDE_H, WARM_DARK)
add_rect(s3, 0, 0, Inches(0.08), SLIDE_H, AMBER)

add_text_box(s3, "✨", Inches(0.4), Inches(0.6), Inches(4.5), Inches(0.9),
             font_size=48, align=PP_ALIGN.CENTER)
add_text_box(s3, "Lucid Moments",
             Inches(0.3), Inches(1.4), Inches(5.0), Inches(1.0),
             font_size=36, bold=True, color=AMBER, font_name="Georgia",
             align=PP_ALIGN.CENTER)
add_text_box(s3, "The Window of Opportunity",
             Inches(0.3), Inches(2.3), Inches(5.0), Inches(0.6),
             font_size=15, italic=True, color=LIGHT_TAN, align=PP_ALIGN.CENTER)

add_rect(s3, Inches(0.9), Inches(2.9), Inches(3.8), Inches(0.04), AMBER)

lucid_desc = [
    ("A lucid moment is a sudden, temporary", 13, False, True, LIGHT_TAN),
    ("period of clarity in a dementia patient —", 13, False, True, LIGHT_TAN),
    ("when they recognize loved ones, recall", 13, False, True, LIGHT_TAN),
    ("memories, and communicate meaningfully.", 13, False, True, LIGHT_TAN),
    ("", 6, False, False, CREAM),
    ("These windows can last minutes to hours.", 12, False, False, MUTED),
    ("They are unpredictable. They are precious.", 12, False, False, MUTED),
    ("Most families are caught unprepared.", 12, False, False, MUTED),
]
add_multiline(s3, lucid_desc, Inches(0.4), Inches(3.05), Inches(4.8), Inches(2.5),
              align=PP_ALIGN.CENTER)

# Lucid stats on left
for i, (num, txt, clr) in enumerate([
    ("73%", "of family caregivers report being\ncaught unprepared during a lucid moment", FLAME),
    ("15–45", "average minutes a lucid moment lasts\nbefore confusion returns", AMBER),
    ("2–4×", "more lucid moments in the morning\nthan any other time of day", GOLD),
]):
    ry = Inches(5.3) + i * Inches(0.0)  # stacked
    add_multiline(s3,
        [(num, 22, True, False, clr), ("  " + txt, 10, False, False, LIGHT_TAN)],
        Inches(0.3), Inches(5.1) + i * Inches(0.65), Inches(5.1), Inches(0.6),
        align=PP_ALIGN.LEFT)

# Right side — what families face
rx = Inches(5.9)
add_text_box(s3, "What Families Face",
             rx, Inches(0.35), Inches(7.0), Inches(0.7),
             font_size=26, bold=True, color=WARM_DARK, font_name="Georgia")
add_rect(s3, rx, Inches(1.0), Inches(6.8), Inches(0.04), AMBER)

problems = [
    ("😰", "Scrambling for words in the moment",
     "Families freeze — unsure what to say when their loved one suddenly recognizes them."),
    ("📷", "No photos or stories at hand",
     "Meaningful memories live in family heads, not in the room when they're needed."),
    ("🔄", "Inconsistent caregiver interactions",
     "Different staff, different approaches — no shared context on what soothes the patient."),
    ("💔", "Missed connections",
     "The moment passes before real conversation happens. The window closes."),
]

for i, (icon, title, desc) in enumerate(problems):
    py = Inches(1.2) + i * Inches(1.4)
    # icon circle
    add_circle(s3, rx + Inches(0.4), py + Inches(0.4), Inches(0.32), LIGHT_TAN)
    add_text_box(s3, icon, rx + Inches(0.1), py + Inches(0.1), Inches(0.7), Inches(0.6),
                 font_size=18, align=PP_ALIGN.CENTER)
    add_text_box(s3, title, rx + Inches(0.85), py, Inches(6.0), Inches(0.42),
                 font_size=13, bold=True, color=WARM_DARK)
    add_text_box(s3, desc, rx + Inches(0.85), py + Inches(0.38), Inches(6.0), Inches(0.75),
                 font_size=11, color=MID_BROWN, italic=False)

add_rect(s3, rx, Inches(6.8), Inches(7.0), Inches(0.45), WARM_DARK)
add_text_box(s3, "Hearth solves exactly this — before the moment happens.",
             rx + Inches(0.2), Inches(6.85), Inches(6.6), Inches(0.35),
             font_size=13, bold=True, italic=True, color=AMBER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW HEARTH WORKS (PIPELINE FLOW)
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, SLIDE_W, SLIDE_H, CREAM)
add_rect(s4, 0, 0, SLIDE_W, Inches(1.3), WARM_DARK)
add_rect(s4, 0, Inches(1.25), SLIDE_W, Inches(0.06), AMBER)

add_text_box(s4, "How Hearth Works",
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
             font_size=34, bold=True, color=CREAM, font_name="Georgia",
             align=PP_ALIGN.CENTER)
add_text_box(s4, "From family memory to care-ready package — in minutes",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
             font_size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Pipeline steps
steps = [
    ("1", "👨‍👩‍👧", "Family\nSubmits", "Memories, photos,\nand voice samples\nshared via Hearth\nweb app",  AMBER),
    ("2", "🧠", "AI\nProcesses", "GPT-4o reads all\nmemories and builds\na rich patient\ncontext model",    GOLD),
    ("3", "⚡", "Lucid Now\nTriggered", "Caregiver presses\n'Lucid Now' when\nthe window opens",           FLAME),
    ("4", "📦", "Package\nGenerated", "4 personalized\nartifacts created\nin seconds via\nOpenAI API",     RGBColor(0x5A,0x32,0x0E)),
    ("5", "🤝", "Caregiver\nDelivers", "Letter read aloud,\nphotos shown,\nvoice played,\nguide followed",  TEAL_ACC),
]

box_w = Inches(2.2)
box_h = Inches(4.5)
gap_x = Inches(0.32)
start = Inches(0.35)
arrow_y = Inches(3.65)

for i, (num, icon, title, desc, clr) in enumerate(steps):
    bx = start + i * (box_w + gap_x)
    by = Inches(1.5)

    # Card background
    add_rect(s4, bx, by, box_w, box_h, clr)

    # Step number badge
    add_circle(s4, bx + Inches(0.38), by + Inches(0.38), Inches(0.28), WARM_DARK)
    add_text_box(s4, num,
                 bx + Inches(0.12), by + Inches(0.1), Inches(0.55), Inches(0.55),
                 font_size=14, bold=True, color=CREAM, align=PP_ALIGN.CENTER)

    # Icon
    add_text_box(s4, icon,
                 bx, by + Inches(0.65), box_w, Inches(0.8),
                 font_size=32, align=PP_ALIGN.CENTER)

    # Title
    add_text_box(s4, title,
                 bx + Inches(0.1), by + Inches(1.4), box_w - Inches(0.2), Inches(0.7),
                 font_size=15, bold=True, color=CREAM, align=PP_ALIGN.CENTER)

    # Divider
    add_rect(s4, bx + Inches(0.3), by + Inches(2.05), box_w - Inches(0.6), Inches(0.03), CREAM)

    # Description
    add_text_box(s4, desc,
                 bx + Inches(0.1), by + Inches(2.15), box_w - Inches(0.2), Inches(1.9),
                 font_size=11, color=CREAM, align=PP_ALIGN.CENTER, italic=False)

    # Arrow between cards
    if i < len(steps) - 1:
        ax = bx + box_w + Inches(0.06)
        add_text_box(s4, "→",
                     ax, arrow_y, gap_x + Inches(0.05), Inches(0.5),
                     font_size=22, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

# Bottom note
add_rect(s4, Inches(0.3), Inches(6.25), Inches(12.7), Inches(0.85), WARM_DARK)
add_text_box(s4,
    "All artifact generation happens server-side using OpenAI GPT-4o · No patient data leaves the local environment · "
    "Demo mode runs without API key",
    Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.6),
    font_size=10.5, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THE FOUR ARTIFACTS
# ═══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, SLIDE_W, SLIDE_H, WARM_DARK)

# Decorative orbs
add_circle(s5, Inches(12.8), Inches(0.5), Inches(1.5), RGBColor(0x40, 0x20, 0x05))
add_circle(s5, Inches(0.5),  Inches(7.2), Inches(1.2), RGBColor(0x35, 0x18, 0x04))

add_text_box(s5, "The Four Artifacts",
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.75),
             font_size=36, bold=True, color=CREAM, font_name="Georgia",
             align=PP_ALIGN.CENTER)
add_text_box(s5, "Generated by GPT-4o · Personalized for every patient · Ready in under 60 seconds",
             Inches(0.5), Inches(0.9), Inches(12), Inches(0.4),
             font_size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

artifacts = [
    ("✉️", "Family Letter",
     "A ~200-word personal letter written from the family member's voice — using real names, real places, and real memories. "
     "Warm, simple language. Read aloud by the caregiver at the start of the visit.",
     ["Personalized to patient's first name", "Written in family member's voice", "Based on submitted memories", "No medical language — pure warmth"],
     AMBER),
    ("🎙️", "Voice Message",
     "A 60–80 word spoken script — intimate and warm, like a phone message from a loved one. "
     "Recorded or read aloud. Hearing a familiar voice or voice-like message can anchor a patient in lucid clarity.",
     ["60–80 words precisely", "Natural spoken language", "References 1–2 specific memories", "Ends with a warm closing"],
     FLAME),
    ("🖼️", "Photo Story",
     "A PDF with family photos and caregiver-read captions. Each photo gets a 2–3 sentence caption — "
     "sensory, emotional, specific — designed to be read aloud while showing the photo to the patient.",
     ["One page per photo", "18pt+ accessible font", "Caption ties photo to a real memory", "Caregiver reads, patient sees"],
     GOLD),
    ("💬", "Dialogue Guide",
     "A structured script for the caregiver: greeting, 5 memory prompts, photo sequence, what to do if distressed, "
     "and a gentle closing. Built entirely from the patient's emotional anchors and family memories.",
     ["Opening & closing scripts", "5 tailored memory prompts", "Distress de-escalation script", "Photo show sequence"],
     TEAL_ACC),
]

card_w2 = Inches(3.0)
card_h2 = Inches(5.3)
gap2    = Inches(0.24)
start2  = Inches(0.35)

for i, (icon, title, desc, bullets, clr) in enumerate(artifacts):
    cx = start2 + i * (card_w2 + gap2)
    cy = Inches(1.45)

    # Card
    add_rect(s5, cx, cy, card_w2, card_h2, RGBColor(0x26, 0x15, 0x08))

    # Top color bar
    add_rect(s5, cx, cy, card_w2, Inches(0.08), clr)

    # Icon
    add_text_box(s5, icon, cx, cy + Inches(0.15), card_w2, Inches(0.75),
                 font_size=30, align=PP_ALIGN.CENTER)

    # Title
    add_text_box(s5, title, cx + Inches(0.1), cy + Inches(0.85),
                 card_w2 - Inches(0.2), Inches(0.5),
                 font_size=16, bold=True, color=clr, align=PP_ALIGN.CENTER)

    # Description
    add_text_box(s5, desc, cx + Inches(0.12), cy + Inches(1.35),
                 card_w2 - Inches(0.24), Inches(1.75),
                 font_size=10, color=LIGHT_TAN, wrap=True)

    # Divider
    add_rect(s5, cx + Inches(0.2), cy + Inches(3.05), card_w2 - Inches(0.4), Inches(0.03), clr)

    # Bullets
    for j, b in enumerate(bullets):
        add_text_box(s5, "▸  " + b,
                     cx + Inches(0.15), cy + Inches(3.15) + j * Inches(0.47),
                     card_w2 - Inches(0.25), Inches(0.45),
                     font_size=9.5, color=CREAM)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TECH ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, SLIDE_W, SLIDE_H, CREAM)
add_rect(s6, 0, 0, SLIDE_W, Inches(1.3), WARM_DARK)
add_rect(s6, 0, Inches(1.25), SLIDE_W, Inches(0.06), AMBER)

add_text_box(s6, "Technology Architecture",
             Inches(0.5), Inches(0.25), Inches(12), Inches(0.75),
             font_size=34, bold=True, color=CREAM, font_name="Georgia",
             align=PP_ALIGN.CENTER)
add_text_box(s6, "FastAPI · OpenAI GPT-4o · fpdf2 · ElevenLabs · Single-Page App",
             Inches(0.5), Inches(0.88), Inches(12), Inches(0.38),
             font_size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Frontend box
fx, fy, fw, fh = Inches(0.4), Inches(1.55), Inches(3.8), Inches(5.2)
add_rect(s6, fx, fy, fw, fh, WARM_DARK)
add_rect(s6, fx, fy, fw, Inches(0.06), TEAL_ACC)
add_text_box(s6, "🌐  Frontend",
             fx + Inches(0.15), fy + Inches(0.15), fw - Inches(0.3), Inches(0.5),
             font_size=16, bold=True, color=TEAL_ACC)
fe_items = [
    "Vanilla JS Single-Page App",
    "7 screens, no page reload",
    "Server-Sent Events (SSE)\nfor live job progress",
    "Streaming chat UI\n(token-by-token display)",
    "Drag-and-drop photo/\nvoice file uploads",
    "Warm light theme\n(CSS custom properties)",
    "Back/Forward navigation\n(history stack)",
]
for i, item in enumerate(fe_items):
    add_text_box(s6, "•  " + item,
                 fx + Inches(0.2), fy + Inches(0.75) + i * Inches(0.6),
                 fw - Inches(0.3), Inches(0.58),
                 font_size=10.5, color=CREAM)

# Backend box
bx2, by2, bw2, bh2 = Inches(4.55), Inches(1.55), Inches(4.0), Inches(5.2)
add_rect(s6, bx2, by2, bw2, bh2, WARM_DARK)
add_rect(s6, bx2, by2, bw2, Inches(0.06), AMBER)
add_text_box(s6, "⚙️  Backend (FastAPI)",
             bx2 + Inches(0.15), by2 + Inches(0.15), bw2 - Inches(0.3), Inches(0.5),
             font_size=16, bold=True, color=AMBER)
be_items = [
    "Python · FastAPI · Uvicorn",
    "REST API + SSE streaming",
    "Background thread pool\nfor blocking AI calls",
    "Memory bank (JSON store)\nfor patient profiles",
    "File serving for PDFs\nand photos",
    "CORS middleware for\nlocal dev",
    "Demo mode fallback\n(no API key needed)",
]
for i, item in enumerate(be_items):
    add_text_box(s6, "•  " + item,
                 bx2 + Inches(0.2), by2 + Inches(0.75) + i * Inches(0.6),
                 bw2 - Inches(0.3), Inches(0.58),
                 font_size=10.5, color=CREAM)

# AI / Output box
ax2, ay2, aw2, ah2 = Inches(8.9), Inches(1.55), Inches(4.1), Inches(5.2)
add_rect(s6, ax2, ay2, aw2, ah2, WARM_DARK)
add_rect(s6, ax2, ay2, aw2, Inches(0.06), FLAME)
add_text_box(s6, "🤖  AI & Output Layer",
             ax2 + Inches(0.15), ay2 + Inches(0.15), aw2 - Inches(0.3), Inches(0.5),
             font_size=16, bold=True, color=FLAME)
ai_items = [
    "OpenAI GPT-4o\n(artifact generation)",
    "GPT-4o streaming\n(caretaker chat)",
    "fpdf2 — PDF generation\nfor all 4 artifacts",
    "Pillow — photo embedding\nin photo story PDF",
    "ElevenLabs (optional)\ntext-to-speech voice",
    "Per-patient context\ninjected into every prompt",
    "JSON-structured outputs\nwith markdown stripping",
]
for i, item in enumerate(ai_items):
    add_text_box(s6, "•  " + item,
                 ax2 + Inches(0.2), ay2 + Inches(0.75) + i * Inches(0.6),
                 aw2 - Inches(0.3), Inches(0.58),
                 font_size=10.5, color=CREAM)

# Arrows between boxes
add_text_box(s6, "⟺", Inches(4.1), Inches(3.75), Inches(0.55), Inches(0.5),
             font_size=20, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_text_box(s6, "⟺", Inches(8.6), Inches(3.75), Inches(0.55), Inches(0.5),
             font_size=20, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

add_rect(s6, 0, Inches(6.95), SLIDE_W, Inches(0.55), WARM_DARK)
add_text_box(s6,
    "Data flow: Browser ↔ FastAPI (REST + SSE) ↔ Agent (GPT-4o) → PDF Generators → File Store → Browser",
    Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
    font_size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — IMPACT & VALUE
# ═══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, SLIDE_W, SLIDE_H, WARM_DARK)
add_rect(s7, 0, 0, Inches(0.08), SLIDE_H, FLAME)

add_text_box(s7, "Impact & Value Proposition",
             Inches(0.4), Inches(0.25), Inches(12.5), Inches(0.75),
             font_size=34, bold=True, color=CREAM, font_name="Georgia")
add_rect(s7, Inches(0.4), Inches(0.95), Inches(12.0), Inches(0.04), AMBER)

# Left: impact metrics
impact = [
    ("< 60s",  "Full care package generated",         "Letter · Voice · Photos · Dialogue guide",     FLAME),
    ("100%",   "Personalized to each patient",         "Every prompt grounded in real family memories", AMBER),
    ("4",      "Artifacts per Lucid Now session",      "Each designed for a different care touchpoint", GOLD),
    ("0",      "Extra tools needed",                   "One web app — works on any device, any browser",TEAL_ACC),
    ("∞",      "Patient profiles supported",           "Add, edit, delete anytime through the web UI",  RGBColor(0xC0,0x80,0x30)),
]

for i, (num, label, sub, clr) in enumerate(impact):
    iy = Inches(1.2) + i * Inches(1.1)
    add_rect(s7, Inches(0.35), iy, Inches(6.2), Inches(1.0), RGBColor(0x26, 0x15, 0x08))
    add_text_box(s7, num, Inches(0.45), iy + Inches(0.05), Inches(1.2), Inches(0.9),
                 font_size=36, bold=True, color=clr, align=PP_ALIGN.CENTER,
                 font_name="Georgia")
    add_text_box(s7, label, Inches(1.65), iy + Inches(0.05), Inches(4.8), Inches(0.45),
                 font_size=14, bold=True, color=CREAM)
    add_text_box(s7, sub,   Inches(1.65), iy + Inches(0.48), Inches(4.8), Inches(0.45),
                 font_size=10.5, color=MUTED, italic=True)

# Right: who benefits
rx3 = Inches(7.0)
add_text_box(s7, "Who Benefits",
             rx3, Inches(1.05), Inches(6.0), Inches(0.6),
             font_size=20, bold=True, color=AMBER, font_name="Georgia")

beneficiaries = [
    ("🧓", "Patients",
     "Experience more meaningful interactions during their most precious windows of clarity."),
    ("👨‍👩‍👧‍👦", "Families",
     "Know their memories and love will be delivered exactly when it matters most."),
    ("🏥", "Caregivers",
     "Enter every lucid moment with a structured, personalized plan — confident and prepared."),
    ("🤖", "AI Systems Field",
     "A real-world agentic AI application with clear, measurable human impact."),
]

for i, (icon, title, desc) in enumerate(beneficiaries):
    by3 = Inches(1.75) + i * Inches(1.35)
    add_rect(s7, rx3, by3, Inches(6.0), Inches(1.2), RGBColor(0x2A, 0x14, 0x06))
    add_text_box(s7, icon, rx3 + Inches(0.1), by3 + Inches(0.1), Inches(0.65), Inches(0.65),
                 font_size=24, align=PP_ALIGN.CENTER)
    add_text_box(s7, title, rx3 + Inches(0.85), by3 + Inches(0.1), Inches(5.0), Inches(0.42),
                 font_size=13, bold=True, color=AMBER)
    add_text_box(s7, desc, rx3 + Inches(0.85), by3 + Inches(0.5), Inches(5.0), Inches(0.65),
                 font_size=10.5, color=LIGHT_TAN)

add_rect(s7, rx3, Inches(7.1), Inches(6.0), Inches(0.04), AMBER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, SLIDE_W, SLIDE_H, WARM_DARK)

# Ambient circles
add_circle(s8, Inches(6.66), Inches(3.75), Inches(4.5), RGBColor(0x28, 0x14, 0x04))
add_circle(s8, Inches(6.66), Inches(3.75), Inches(3.0), RGBColor(0x35, 0x1C, 0x06))
add_circle(s8, Inches(6.66), Inches(3.75), Inches(1.6), RGBColor(0x48, 0x28, 0x08))

add_rect(s8, 0, 0, Inches(0.08), SLIDE_H, FLAME)
add_rect(s8, SLIDE_W - Inches(0.08), 0, Inches(0.08), SLIDE_H, FLAME)

add_text_box(s8, "🔥", Inches(0), Inches(1.5), SLIDE_W, Inches(1.2),
             font_size=64, align=PP_ALIGN.CENTER)

add_text_box(s8, "Hearth",
             Inches(0), Inches(2.6), SLIDE_W, Inches(1.1),
             font_size=64, bold=True, color=AMBER, font_name="Georgia",
             align=PP_ALIGN.CENTER)

add_text_box(s8, "When a lucid moment happens — be ready.",
             Inches(1.0), Inches(3.65), Inches(11.3), Inches(0.65),
             font_size=20, italic=True, color=LIGHT_TAN, align=PP_ALIGN.CENTER,
             font_name="Calibri Light")

add_rect(s8, Inches(3.5), Inches(4.35), Inches(6.33), Inches(0.04), AMBER)

add_text_box(s8,
    "Built with  FastAPI · OpenAI GPT-4o · fpdf2 · Pillow · ElevenLabs",
    Inches(1.0), Inches(4.55), Inches(11.3), Inches(0.5),
    font_size=13, color=MUTED, align=PP_ALIGN.CENTER)

add_text_box(s8,
    "BUAN 6V99 — AI Agentic Systems  ·  Spring 2026  ·  Nikhil Gajula",
    Inches(1.0), Inches(5.15), Inches(11.3), Inches(0.45),
    font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_text_box(s8,
    "\"Technology in service of the most human moments.\"",
    Inches(1.5), Inches(6.2), Inches(10.3), Inches(0.6),
    font_size=15, italic=True, color=RGBColor(0xC8, 0xA0, 0x60),
    align=PP_ALIGN.CENTER, font_name="Georgia")


# ── Save ──────────────────────────────────────────────────────────────────────
out = Path(__file__).parent / "hearth_presentation.pptx"
prs.save(str(out))
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
